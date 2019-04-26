#!/usr/bin/env python3
import pprint
from flask import Flask, render_template,request,flash, url_for, request, session, redirect
from flask_pymongo import PyMongo
import bcrypt
import numpy as np
import pandas as pd
from itertools import islice
# from sklearn.preprocessing import label_binarize
# from sklearn.preprocessing import LabelEncoder
# from sklearn.linear_model import LogisticRegression
# from sklearn.metrics import precision_score
# from sklearn.metrics import recall_score
# from sklearn.model_selection import cross_val_score
# from sklearn.ensemble import AdaBoostClassifier
import math
from pptx import Presentation
import statistics
from datetime import date
from werkzeug.security import generate_password_hash, check_password_hash
from fractions import Fraction


app = Flask(__name__)

app.config['SECRET_KEY'] = "a95ec5d5d69e5e75d869feb78c35d2e15090ae5e2259a5b6"
app.config["MONGO_URI"] = "mongodb://localhost:27017/apdc"

mongo = PyMongo(app)
pp = pprint.PrettyPrinter(indent=4)

score = 0

def playerScore(level, exp):
    totalScore = 0

    if level == 'Beginner':
        totalScore += 2
    elif level == 'Intermediate':
        totalScore += 5
    else:
        totalScore += 9

    if exp == 10:
        totalScore += 10
    elif exp == 5:
        totalScore += 5
    else:
        totalScore += 2

    return totalScore

def teamScore(tot_member, totalScore):
    avgValue = totalScore / tot_member
    teamPercent = avgValue / 20 *100
    return teamPercent

def calculate_aspect(width: int, height: int) -> str:
    temp = 0

    def gcd(a, b):
        return a if b == 0 else gcd(b, a % b)

    if width == height:
        return "1:1"

    if width < height:
        temp = width
        width = height
        height = temp

    divisor = gcd(width, height)

    x = int(width / divisor) if not temp else int(height / divisor)
    y = int(height / divisor) if not temp else int(width / divisor)

    return f"{x}:{y}"


def startupRankLoc(location):
    startup_dict = {}
    # csv = open("START.csv","r","encoding = utf-8")
    dataset = pd.read_csv("datasets/startup_ranking.csv",delimiter=',', names = ['rank', 'startup', 'startup_logo', 'sr_score', 'description', 'location','none'])
    startupLoc = dataset[dataset.location == location].head(10)
    startupLoc = startupLoc.to_dict('index')
    # pp.pprint(startupLoc)
    return startupLoc

def calculateROI(invest,tot):
    invest_gained = invest
    total_required_investment = tot
    roi = ((abs(invest_gained - total_required_investment) / total_required_investment) * 100)
    #print('The ROI in % is :',roi)
    #IRR Calculation
    d0 = date(2018, 4, 15)
    d1 = date(2020, 4, 15)
    delta = d1 - d0
    year = math.floor(delta.days / 365)
    roi_year = ((invest_gained - total_required_investment / total_required_investment)**(1/year)-1)
    #print('The Year is:',round(year,2))
    #print('The ROR for',year,'years is:',round(roi_year,2))
    return roi,roi_year

def read_csv(category):
    startup_dict = {}
    # csv = open("START.csv","r","encoding = utf-8")
    dataset = pd.read_csv("datasets/START.csv")
    name = (dataset.iloc[:, 0:1]).values.tolist()
    logo = (dataset.iloc[:, 6:7]).values.tolist()
    status = (dataset.iloc[:, 1:2]).values.tolist()
    fund = (dataset.iloc[:, 8:9]).values.tolist()
    website_link = (dataset.iloc[:, 9:10]).values.tolist()
    founded = (dataset.iloc[:, 2:3]).values.tolist()
    categories = (dataset.iloc[:, 4:5]).values.tolist()
    # print(categories)
    categories = sum(categories, [])
    name = sum(name, [])
    logo = sum(logo, [])
    status = sum(status, [])
    fund = sum(fund,[])
    website_link = sum(website_link, [])
    founded = sum(founded, [])
    startup_list = []
    startup_key = ''
    startup_name = []
    startup_logo = []
    startup_foundedat =[]
    startup_website = []
    startup_status = []

    for j in range(len(categories)-1):
        if (category) in str(categories[j]):
            startup_key = name[j]
            startup_list.append(logo[j])
            startup_list.append(founded[j])
            startup_list.append(website_link[j])
            startup_list.append(status[j])
            # startup_list.append(fund[j])
            startup_dict[startup_key] = startup_list
            startup_list = []
    # print(startup_name, startup_foundedat,)

    return startup_dict

def take(n, iterable):
    "Return first n items of the iterable as a list"
    return list(islice(iterable, n))

@app.route('/')
def index():
    if 'username' in session:
        name = session['username']
        return render_template('index.html' , username = name)

    else:
        return render_template('index.html')

@app.route('/startup_comp', methods = ['GET', 'POST'])
def startup_comparator():
    if request.method == "POST":
        startup = mongo.db.startup
        p_type = request.form['p_type']
        startup_dict = read_csv(p_type)
        startup.insert({'products':[ {'product_name' : request.form['p_name'], 'product_type' : request.form['p_type'], 'usp' : request.form['usp']} ] })
        return render_template('startup_compare.html', dic = startup_dict ,name = True)
    else:
        return render_template('startup_compare.html' ,name = False)

@app.route('/ratio',methods=['GET','POST'])
def cac_ratio():
    if request.method == 'POST':
        clv_cac = mongo.db.users
        login_user = clv_cac.find_one({'name':session['name']})
        #print(login_user)
        tot_acqui = int(request.form['arc'])
        noCust =int (request.form['acl'])
        avgorder = int(request.form['avgorder'])
        noorder = int(request.form['nooforder'])
        uniqueCust = int(request.form['uniqueCust'])
        pf =  int(noorder/uniqueCust)
        pro = int(request.form['profit'])
        cac = int(tot_acqui / noCust)
        clv = int(avgorder * pf * uniqueCust)
        ratio = calculate_aspect(clv,cac)
        print(ratio)
        clv_cac.update_one(
                                {"_id": login_user["_id"]},
                                {"$set":
                                    {'ratio':[
                                                {
                                                    'total_acquistion_cost':tot_acqui,'no_customer':noCust,'average_order':avgorder,'no_order':noorder,'unique_customer':uniqueCust,'purchase_frequency':pf,'profit':pro,'cac':cac,'clv':clv,'ratio':ratio
                                                }
                                             ]
                                    }
                                }
                            )
        return redirect('/investments')
    return render_template('admin/values.html')


@app.route('/register', methods=['POST', 'GET'])
def register():
    if request.method == 'POST':
        users = mongo.db.users
        existing_user = users.find_one({'name' : request.form['name']})
        password = request.form['pass']

        if existing_user is None:
            hashpass = generate_password_hash(password)
            users.insert({'name' : request.form['name'], 'password' : hashpass,'email' :request.form['email'],'contact_number':request.form['phone'],'type':request.form['user_type']})
            session['name'] = request.form['name']
            return redirect('/info')
        return 'That username already exists!'

    return render_template('admin/register.html')

@app.route('/login',methods=['GET','POST'])
def getLogin():
    if request.method == 'POST':
        users = mongo.db.users
        login_user = users.find_one({'name' : request.form['username']})
        #pp.pprint(login_user['type'])
        password = request.form['pass']
        if login_user:
            if check_password_hash(login_user['password'],password):
               session['name'] = request.form['username']
               if login_user['type'] == 'Entrepreneur':
                   return redirect('/dashboard')
            #    print('login ')
        return redirect(url_for('index'))

    return render_template('admin/login.html')

@app.route('/logout')
def logout():
   # remove the username from the session if it is there
   session.pop('username', None)
   session.pop('name', None)
   return redirect(url_for('index'))


@app.route('/info',methods=['GET','POST'])
def info():
    if request.method == 'POST':
        info = mongo.db.users
        login_user = info.find_one({'name':session['name']})
        info.update_one({"_id": login_user["_id"]}, {"$set": {'info':[{'company_name':request.form['name'],'product_info':request.form['product_info'],'product_type':request.form['product_type'],'product_base':request.form['product_base'],'product_name':request.form['product_name'],'usp':request.form['usp'],'location':request.form['location']}]} })
        return redirect('/ratio')
    else:
        return render_template('admin/info.html')

@app.route('/dashboard',methods = ['GET','POST'])
def dash():
    if request.method == 'GET':
        score = 0
        info = mongo.db.users
        login_user = info.find_one({'name':session['name']})

        user_info = login_user['info']
        user_ratio = login_user['ratio']

        team_details = login_user['team_details']
        tot_member = len(team_details)
        prod_type = user_info[0]['product_type']
        location =login_user['info'][0]['location']

        invest_gained = login_user['investment_details'][0]['investment_gained']
        total_required_investment = login_user['investment_details'][0]['required_investment']

        startup_dict  = read_csv(prod_type)
        n_items = take(10, startup_dict.items())
        rankStartups = startupRankLoc(location)
        r_items = take(9, rankStartups.items())

        for i in team_details:
            level = i['level']
            exp = i['experience']
            score += playerScore(level , exp)


        team_percent = int(teamScore(tot_member, score))

        rates = calculateROI(int(invest_gained),int(total_required_investment))
        total_returns = list(rates)
        teamPercent = team_percent
        total_returns.append(teamPercent)
        riskFactor = statistics.mean(total_returns)

        pp.pprint(login_user)
        return render_template('admin/dashboard.html', user_info = user_info,
        startup_dict = startup_dict,
        login_user = login_user,
        startups = n_items,
        locstartups = r_items,
        location = location,
        team_details= team_details,
        tot_member=tot_member,
        team_percent = team_percent,
        riskfact = round(riskFactor,2),
        user_ratio = user_ratio
        )
    else:
        info = mongo.db.users
        login_user = info.find_one({'name':session['name']})
        pp.pprint(login_user)
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = login_user['info'][0]['product_name']
        subtitle.text = login_user['info'][0]['company_name']
        prs.save('test.pptx')




@app.route('/samestartups')
def similarStartups():
    info = mongo.db.users
    login_user = info.find_one({'name':session['name']})
    user_info = login_user['info']
    prod_type = user_info[0]['product_type']
    startup_dict  = read_csv(prod_type)
    # pp.pprint(startup_dict)
    return render_template('admin/startup_compare.html', user_info = user_info, startup_dict= startup_dict, login_user = login_user)

@app.route('/addinvestments',methods=['GET','POST'])
def addInvestments():
    if request.method == 'POST':
        pass
    else:
        info = mongo.db.users
        login_user = info.find_one({'name':session['name']})
        user_info = login_user['info']
        return render_template('admin/investment_details.html', user_info = user_info,login_user = login_user)


@app.route('/add-team',methods=['GET','POST'])
def addTeamProfile():
    if request.method == 'POST':
        team_details = mongo.db.users
        login_user = team_details.find_one({'name':session['name']})
        name = request.form['name']
        position = request.form['position']
        exp = request.form['experience']
        level = request.form['level']
        team_details.update(
                                {"_id": login_user["_id"]},
                                {"$push":
                                    {'team_details':
                                                {
                                                    'name' : name,'positon':position,'experience':exp,'level':level
                                                }

                                    }
                                }
                            )
        return render_template('/admin/teamprofile.html')
    else:
        return render_template('/admin/teamprofile.html')

@app.route('/investments',methods=['GET','POST'])
def addInvestmentDetails():
    if request.method == 'POST':
        investment_details = mongo.db.users
        login_user = investment_details.find_one({'name':session['name']})
        invest_gain = request.form['invesment_gained']
        total_req_invest = request.form['investment_required']
        seed_fund = request.form['seed_fund']
        business_model = request.form['business_model']
        investment_details.update_one(
                                {"_id": login_user["_id"]},
                                {"$set":
                                    {'investment_details':[
                                                {
                                                    'investment_gained':invest_gain,'required_investment':total_req_invest,'seedFunding':seed_fund,'business_model':business_model
                                                }
                                             ]
                                    }
                                }
                            )
        return redirect('/add-team')
    else:
        return render_template('/admin/investment-details.html')

@app.route('/calcRoi', methods = ['GET','POST'])
def calcRoi():
    if request.method == 'POST':
        investment_details = mongo.db.users
        login_user = investment_details.find_one({'name':session['name']})
        pp.pprint(login_user['investment_details'][0]['required_investment'])
        invest_gained = login_user['investment_details'][0]['investment_gained']
        total_required_investment = login_user['investment_details'][0]['required_investment']
        roi = float((abs(invest_gained - total_required_investment) / total_required_investment) * 100)
        #print('The ROI in % is :',roi)
        #IRR Calculation
        d0 = date(request.form['date1'])
        d1 = date(request.form['date2'])
        delta = d1 - d0
        year = math.floor(delta.days / 365)
        roi_year = ((invest_gained - total_required_investment / total_required_investment)**(1/year)-1)
        #print('The Year is:',round(year,2))
        #print('The ROR for',year,'years is:',round(roi_year,2))
        return roi,roi_year
    return render_template('/admin/guage.html')

if __name__ == '__main__':
    app.run(host='0.0.0.0',debug=True)
